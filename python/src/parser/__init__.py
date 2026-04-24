"""Document text extraction helpers."""

from __future__ import annotations

import csv
import json
from pathlib import Path
from xml.etree import ElementTree as ET

from src.parser.extractor import DocumentContent, PageContent, extract_pages, extract_text

SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".doc": "Word",
    ".txt": "Text",
    ".md": "Markdown",
    ".log": "Log",
    ".html": "HTML",
    ".htm": "HTML",
    ".epub": "EPUB",
    ".rtf": "RTF",
    ".tex": "LaTeX",
    ".csv": "CSV",
    ".pptx": "PowerPoint",
    ".xlsx": "Excel",
    ".srt": "Subtitle",
    ".json": "JSON",
    ".xml": "XML",
}


def _single_page(path: Path, text: str) -> DocumentContent:
    return DocumentContent(
        pages=[
            PageContent(
                page_num=1,
                text=text,
                width=0,
                height=0,
                is_dual_column=False,
            )
        ],
        source_path=str(path),
    )


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(errors="replace")


def _extract_html(path: Path) -> str:
    raw = _read_text_file(path)
    try:
        from bs4 import BeautifulSoup

        return BeautifulSoup(raw, "html.parser").get_text("\n")
    except ImportError:
        import re

        return re.sub(r"<[^>]+>", " ", raw)


def _extract_rtf(path: Path) -> str:
    raw = _read_text_file(path)
    try:
        from striprtf.striprtf import rtf_to_text

        return rtf_to_text(raw)
    except ImportError:
        return raw


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError("python-docx is required to parse .docx files") from exc

    document = docx.Document(str(path))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append("\t".join(cells))
    return "\n\n".join(paragraphs)


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to parse .xlsx files") from exc

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"# {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in values):
                parts.append("\t".join(values))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to parse .pptx files") from exc

    presentation = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        slide_parts = [f"# Slide {idx}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)


def _extract_csv(path: Path) -> str:
    text = _read_text_file(path)
    rows = list(csv.reader(text.splitlines()))
    return "\n".join("\t".join(row) for row in rows)


def _extract_json(path: Path) -> str:
    data = json.loads(_read_text_file(path))
    return json.dumps(data, ensure_ascii=False, indent=2)


def _extract_xml(path: Path) -> str:
    raw = _read_text_file(path)
    try:
        root = ET.fromstring(raw)
        parts: list[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
        return "\n".join(parts) or raw
    except ET.ParseError:
        return raw


def _extract_epub(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
        from ebooklib import ITEM_DOCUMENT, epub
    except ImportError as exc:
        raise RuntimeError("EbookLib and beautifulsoup4 are required to parse .epub files") from exc

    book = epub.read_epub(str(path))
    parts: list[str] = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        text = soup.get_text("\n").strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def extract_document(path: str | Path) -> DocumentContent:
    file_path = Path(path)
    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file format: {ext}")
    if ext == ".pdf":
        return extract_pages(file_path)
    if ext == ".doc":
        raise RuntimeError(".doc files are not supported directly. Please save as .docx first.")

    if ext == ".docx":
        text = _extract_docx(file_path)
    elif ext == ".xlsx":
        text = _extract_xlsx(file_path)
    elif ext == ".pptx":
        text = _extract_pptx(file_path)
    elif ext in {".html", ".htm"}:
        text = _extract_html(file_path)
    elif ext == ".rtf":
        text = _extract_rtf(file_path)
    elif ext == ".csv":
        text = _extract_csv(file_path)
    elif ext == ".json":
        text = _extract_json(file_path)
    elif ext == ".xml":
        text = _extract_xml(file_path)
    elif ext == ".epub":
        text = _extract_epub(file_path)
    else:
        text = _read_text_file(file_path)

    return _single_page(file_path, text)


__all__ = ["extract_text", "extract_pages", "extract_document", "SUPPORTED_EXTENSIONS", "DocumentContent"]
