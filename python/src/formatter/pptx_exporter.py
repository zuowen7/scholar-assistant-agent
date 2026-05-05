"""论文转 PPTX 导出 — 借鉴 nature-paper2ppt 的叙事弧线设计。

核心原则：
- 以论文的科学论证为叙事主线，而非机械按章节顺序
- 每张幻灯片承载一个论点
- 图片优先于文字，hero figure + 窄注释栏
- 默认中文学术报告风格，保留关键技术术语英文

输出：16:9 宽屏 .pptx，中文标题 + 演讲者备注。
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ── 设计常量 ──────────────────────────────────────────────────────────────────

SLIDE_WIDTH = Inches(13.333)   # 16:9
SLIDE_HEIGHT = Inches(7.5)

# Nature 风格配色
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x27, 0x27, 0x27)
COLOR_ACCENT = RGBColor(0x0F, 0x4D, 0x92)       # Nature blue
COLOR_SECONDARY = RGBColor(0x76, 0x76, 0x76)    # 灰色
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_BORDER = RGBColor(0xD8, 0xD8, 0xD8)
COLOR_GAIN = RGBColor(0x2E, 0x9E, 0x44)
COLOR_LOSS = RGBColor(0xE5, 0x39, 0x35)

# 字体（优先 Arial，不可用时回退系统默认）
FONT_TITLE = "Arial"
FONT_BODY = "Arial"
FONT_CJK = None  # python-pptx 会自动使用系统 CJK 字体

# ── 幻灯片内容结构 ────────────────────────────────────────────────────────────

SLIDE_TEMPLATES = {
    "title": {
        "title_size": Pt(36),
        "subtitle_size": Pt(18),
        "title_color": COLOR_ACCENT,
        "subtitle_color": COLOR_SECONDARY,
    },
    "section": {
        "title_size": Pt(28),
        "body_size": Pt(16),
        "title_color": COLOR_ACCENT,
        "body_color": COLOR_TEXT,
    },
    "content": {
        "title_size": Pt(24),
        "body_size": Pt(14),
        "title_color": COLOR_ACCENT,
        "body_color": COLOR_TEXT,
    },
    "result": {
        "title_size": Pt(22),
        "body_size": Pt(13),
        "title_color": COLOR_ACCENT,
        "body_color": COLOR_TEXT,
    },
    "conclusion": {
        "title_size": Pt(24),
        "body_size": Pt(16),
        "title_color": COLOR_TEXT,
        "body_color": COLOR_SECONDARY,
    },
}


class PPTXBuilder:
    """构建学术报告 PPTX 的 Builder"""

    def __init__(self):
        if not HAS_PPTX:
            raise ImportError(
                "python-pptx 未安装。请运行: pip install python-pptx"
            )
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.slide_count = 0

    def _add_blank_slide(self) -> Any:
        """添加空白幻灯片"""
        layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(layout)
        self.slide_count += 1
        return slide

    def _set_slide_bg(self, slide: Any, color: RGBColor = COLOR_BG) -> None:
        """设置幻灯片背景色"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(
        self,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: Pt = Pt(14),
        color: RGBColor = COLOR_TEXT,
        bold: bool = False,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        font_name: str = FONT_BODY,
    ) -> Any:
        """添加文本框"""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_bullet_list(
        self,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        items: list[str],
        font_size: Pt = Pt(14),
        color: RGBColor = COLOR_TEXT,
    ) -> Any:
        """添加项目符号列表"""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = font_size
            p.font.color.rgb = color
            p.font.name = FONT_BODY
            p.level = 0
            p.space_after = Pt(6)
            # 添加项目符号
            pPr = p._pPr
            if pPr is None:
                from pptx.oxml.ns import qn
                pPr = p._p.get_or_add_pPr()
            from pptx.oxml.ns import qn
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
            # 移除旧的项目符号设置
            for child in list(pPr):
                if child.tag in (qn('a:buChar'), qn('a:buNone')):
                    pPr.remove(child)
            pPr.append(buChar)

        return txBox

    def _add_speaker_notes(self, slide: Any, notes_text: str) -> None:
        """添加演讲者备注"""
        if not notes_text:
            return
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    def _add_source_label(
        self, slide: Any, source_text: str
    ) -> None:
        """在幻灯片底部添加来源标签"""
        self._add_textbox(
            slide,
            left=0.5, top=7.0, width=12, height=0.3,
            text=source_text,
            font_size=Pt(8),
            color=COLOR_SECONDARY,
        )

    # ── 幻灯片构建方法 ──────────────────────────────────────────────────────

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        authors: str = "",
        venue: str = "",
        notes: str = "",
    ) -> None:
        """标题页"""
        slide = self._add_blank_slide()

        # 顶部装饰线
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(1), Inches(2.0), Inches(11.333), Pt(3),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.fill.background()

        # 标题
        self._add_textbox(
            slide, left=1, top=2.2, width=11.333, height=1.5,
            text=title,
            font_size=SLIDE_TEMPLATES["title"]["title_size"],
            color=SLIDE_TEMPLATES["title"]["title_color"],
            bold=True,
        )

        # 副标题
        if subtitle:
            self._add_textbox(
                slide, left=1, top=3.5, width=11.333, height=0.8,
                text=subtitle,
                font_size=SLIDE_TEMPLATES["title"]["subtitle_size"],
                color=SLIDE_TEMPLATES["title"]["subtitle_color"],
            )

        # 作者 / 期刊
        if authors:
            self._add_textbox(
                slide, left=1, top=5.0, width=11.333, height=0.5,
                text=authors,
                font_size=Pt(12),
                color=COLOR_SECONDARY,
            )
        if venue:
            self._add_textbox(
                slide, left=1, top=5.4, width=11.333, height=0.4,
                text=venue,
                font_size=Pt(10),
                color=COLOR_SECONDARY,
            )

        self._add_speaker_notes(slide, notes)

    def add_section_slide(
        self,
        title: str,
        bullets: list[str],
        notes: str = "",
        source: str = "",
    ) -> None:
        """章节/内容幻灯片"""
        slide = self._add_blank_slide()

        # 标题
        self._add_textbox(
            slide, left=0.8, top=0.5, width=11.733, height=0.8,
            text=title,
            font_size=SLIDE_TEMPLATES["content"]["title_size"],
            color=SLIDE_TEMPLATES["content"]["title_color"],
            bold=True,
        )

        # 标题下划线
        line = slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.2), Inches(2), Pt(2),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BORDER
        line.line.fill.background()

        # 正文要点
        if bullets:
            self._add_bullet_list(
                slide, left=0.8, top=1.6, width=11.733, height=5.0,
                items=bullets,
                font_size=SLIDE_TEMPLATES["content"]["body_size"],
            )

        if source:
            self._add_source_label(slide, source)
        self._add_speaker_notes(slide, notes)

    def add_result_slide(
        self,
        title: str,
        bullets: list[str],
        figure_path: str = "",
        notes: str = "",
        source: str = "",
    ) -> None:
        """结果幻灯片（图为主 + 窄注释栏）"""
        slide = self._add_blank_slide()

        # 标题
        self._add_textbox(
            slide, left=0.5, top=0.3, width=12.333, height=0.7,
            text=title,
            font_size=SLIDE_TEMPLATES["result"]["title_size"],
            color=SLIDE_TEMPLATES["result"]["title_color"],
            bold=True,
        )

        # 如果有图，采用不对称布局：图占左侧 70%，要点占右侧 25%
        if figure_path and Path(figure_path).exists():
            try:
                # 图片放在左侧
                pic = slide.shapes.add_picture(
                    str(figure_path),
                    Inches(0.5), Inches(1.2),
                    Inches(8.5), Inches(5.5),
                )
                # 要点放在右侧窄栏
                if bullets:
                    self._add_bullet_list(
                        slide, left=9.3, top=1.5, width=3.7, height=5.0,
                        items=bullets,
                        font_size=Pt(11),
                        color=COLOR_TEXT,
                    )
            except Exception as e:
                logger.warning("插入图片失败 %s: %s", figure_path, e)
                if bullets:
                    self._add_bullet_list(
                        slide, left=0.5, top=1.5, width=12.333, height=5.0,
                        items=bullets,
                        font_size=SLIDE_TEMPLATES["result"]["body_size"],
                    )
        else:
            # 无图时使用全宽要点
            if bullets:
                self._add_bullet_list(
                    slide, left=0.5, top=1.2, width=12.333, height=5.5,
                    items=bullets,
                    font_size=SLIDE_TEMPLATES["result"]["body_size"],
                )

        if source:
            self._add_source_label(slide, source)
        self._add_speaker_notes(slide, notes)

    def add_conclusion_slide(
        self,
        title: str,
        key_points: list[str],
        limitations: list[str] | None = None,
        notes: str = "",
    ) -> None:
        """结论幻灯片"""
        slide = self._add_blank_slide()

        # 标题
        self._add_textbox(
            slide, left=0.8, top=0.5, width=11.733, height=0.8,
            text=title,
            font_size=SLIDE_TEMPLATES["conclusion"]["title_size"],
            color=SLIDE_TEMPLATES["conclusion"]["title_color"],
            bold=True,
        )

        # 核心要点
        if key_points:
            self._add_bullet_list(
                slide, left=0.8, top=1.5, width=11.733, height=3.5,
                items=key_points,
                font_size=SLIDE_TEMPLATES["conclusion"]["body_size"],
                color=COLOR_TEXT,
            )

        # 局限性（使用灰色小字）
        if limitations:
            limit_text = "局限性: " + "; ".join(limitations)
            self._add_textbox(
                slide, left=0.8, top=5.5, width=11.733, height=1.0,
                text=limit_text,
                font_size=Pt(11),
                color=COLOR_SECONDARY,
            )

        self._add_speaker_notes(slide, notes)

    def save(self, output_path: str | Path) -> Path:
        """保存 PPTX 文件"""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info("PPTX 已保存: %s (%d slides)", output_path, self.slide_count)
        return output_path


# ── 便捷函数 ──────────────────────────────────────────────────────────────────

def build_paper_pptx(
    output_path: str | Path,
    *,
    paper_title: str = "",
    authors: str = "",
    venue: str = "",
    slides: list[dict] | None = None,
) -> Path:
    """从结构化数据构建论文 PPTX。

    Args:
        output_path: 输出文件路径
        paper_title: 论文标题
        authors: 作者列表
        venue: 期刊/会议名
        slides: 幻灯片列表，每项为 dict:
            - type: title / section / result / conclusion
            - title: 幻灯片标题
            - bullets: 要点列表 (可选)
            - figure_path: 图片路径 (可选，result 类型)
            - notes: 演讲者备注 (可选)
            - source: 来源标注 (可选)

    Returns:
        输出文件路径
    """
    builder = PPTXBuilder()

    # 标题页
    builder.add_title_slide(
        title=paper_title or "学术报告",
        subtitle="",
        authors=authors,
        venue=venue,
    )

    # 内容页
    for sd in (slides or []):
        stype = sd.get("type", "section")
        title = sd.get("title", "")
        bullets = sd.get("bullets", [])
        figure_path = sd.get("figure_path", "")
        notes = sd.get("notes", "")
        source = sd.get("source", "")

        if stype == "result":
            builder.add_result_slide(
                title=title,
                bullets=bullets,
                figure_path=figure_path,
                notes=notes,
                source=source,
            )
        elif stype == "conclusion":
            builder.add_conclusion_slide(
                title=title,
                key_points=bullets,
                notes=notes,
            )
        else:  # section / content
            builder.add_section_slide(
                title=title,
                bullets=bullets,
                notes=notes,
                source=source,
            )

    return builder.save(output_path)


def export_translated_paper_to_pptx(
    output_path: str | Path,
    *,
    title: str = "",
    authors: str = "",
    venue: str = "",
    section_slides: list[dict] | None = None,
) -> Path:
    """从翻译结果导出 PPTX。

    与翻译管道集成：将翻译好的章节内容转换为 PPTX 幻灯片。
    每张幻灯片使用 conclusion-style 标题（陈述观点而非仅标注主题）。

    Args:
        output_path: 输出文件路径
        title: 论文标题
        authors: 作者
        venue: 发表期刊/会议
        section_slides: 章节幻灯片数据列表，每项:
            - section_title: 章节标题（如 "Introduction"）
            - section_type: 章节类型
            - key_points: 中文要点列表
            - notes: 演讲者备注

    Returns:
        输出文件路径
    """
    slides_data: list[dict] = []

    for ss in (section_slides or []):
        slides_data.append({
            "type": "section",
            "title": ss.get("section_title", ""),
            "bullets": ss.get("key_points", []),
            "notes": ss.get("notes", ""),
            "source": "",
        })

    return build_paper_pptx(
        output_path,
        paper_title=title,
        authors=authors,
        venue=venue,
        slides=slides_data,
    )
